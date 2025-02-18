from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Stock Price Prediction Report"
subtitle.text = "A Comprehensive Analysis of ARIMA vs Machine Learning Models\nPresented by: [Your Name]"

# Add Slide for Introduction
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "1. Introduction"
content.text = (
    "- Objective: Analyze stock price prediction using ARIMA & ML models.\n"
    "- Scope: Data preparation, EDA, Feature Engineering, Model Development, and Evaluation."
)

# Add Slide for Data Preparation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "2. Data Preparation"
content.text = (
    "- Historical stock data for multiple companies.\n"
    "- Data cleaning: Handling missing values & aligning dates.\n"
    "- Final dataset: `final_merged_stock_data.csv`"
)

# Add Slide for EDA
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "3. Exploratory Data Analysis (EDA)"
content.text = (
    "- Identified stock price trends and volatility.\n"
    "- Tesla (TSLA) exhibited the highest fluctuations.\n"
    "- Strong correlation found between AAPL & MSFT.\n"
    "- Rolling Mean & Price Distribution Analysis."
)

# Add Slide for Feature Engineering
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "4. Feature Engineering"
content.text = (
    "- Created lag features, rolling means, and percentage changes.\n"
    "- Included ATR-based volatility indicators.\n"
    "- Improved model prediction accuracy."
)

# Add Slide for Model Development
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "5. Model Development"
content.text = (
    "Models Used:\n"
    "- ARIMA (Best Order: (1,1,0))\n"
    "- Gradient Boosting\n"
    "- XGBoost\n"
    "- LightGBM (Best performing model)"
)

# Add Slide for Model Evaluation
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "6. Model Evaluation & Comparison"
content.text = (
    "- **ARIMA performed poorly (RMSE: 25.68)**\n"
    "- **Gradient Boosting improved accuracy**\n"
    "- **XGBoost further optimized performance**\n"
    "- **LightGBM achieved the best results (RMSE: 2.810, R²: 0.91)**"
)

# Add Slide for Business Implications
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "7. Business Implications"
content.text = (
    "- Use LightGBM for stock trend forecasting.\n"
    "- Buy/Sell decisions based on model predictions.\n"
    "- Manage risk by avoiding volatile stocks."
)

# Add Slide for Deliverables
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "8. Final Deliverables"
content.text = (
    "- Data Preparation: `final_merged_stock_data.csv`\n"
    "- EDA & Visuals: `eda_analysis.ipynb`\n"
    "- Feature Engineering: `AAPL_boosting_features.csv`\n"
    "- Models & Training Code: `models.py`\n"
    "- Evaluation Report: `model_comparison.csv`\n"
    "- Presentation Slides: `Stock_Prediction_Presentation.pptx`\n"
    "- README Documentation: `README.md`"
)

# Save the PowerPoint file
pptx_filename = "Stock_Prediction_Presentation.pptx"
prs.save(pptx_filename)

# Generate PDF Report
from fpdf import FPDF

class PDF(FPDF):
    def header(self):
        self.set_font("Arial", "B", 16)
        self.cell(0, 10, "Stock Price Prediction Report", ln=True, align="C")
        self.ln(10)

    def chapter_title(self, title):
        self.set_font("Arial", "B", 14)
        self.cell(0, 10, title, ln=True, align="L")
        self.ln(5)

    def chapter_body(self, body):
        self.set_font("Arial", "", 12)
        self.multi_cell(0, 8, body)
        self.ln(5)

pdf = PDF()
pdf.set_auto_page_break(auto=True, margin=15)
pdf.add_page()

pdf.chapter_title("1. Introduction")
pdf.chapter_body(
    "Objective: This report evaluates stock price prediction using ARIMA and "
    "machine learning models. We compare different models to determine the best predictive strategy."
)

pdf.chapter_title("2. Data Preparation")
pdf.chapter_body(
    "Historical stock data was collected and preprocessed to handle missing values, "
    "align dates, and extract relevant features for predictive modeling."
)

pdf.chapter_title("3. Exploratory Data Analysis (EDA)")
pdf.chapter_body(
    "Key findings from EDA include trend identification, volatility analysis, "
    "and correlation between different stock prices."
)

pdf.chapter_title("4. Feature Engineering")
pdf.chapter_body(
    "Created lag features, rolling statistics, and momentum indicators to "
    "enhance predictive performance."
)

pdf.chapter_title("5. Model Development")
pdf.chapter_body(
    "The models evaluated include:\n"
    "- ARIMA (1,1,0)\n"
    "- Gradient Boosting\n"
    "- XGBoost\n"
    "- LightGBM (Best Performing Model)"
)

pdf.chapter_title("6. Model Evaluation & Comparison")
pdf.chapter_body(
    "Performance Metrics:\n"
    "- ARIMA: RMSE: 25.68\n"
    "- Gradient Boosting: RMSE: 2.929\n"
    "- XGBoost: RMSE: 2.810\n"
    "- LightGBM: RMSE: 2.810 (Best Model)"
)

pdf.chapter_title("7. Business Implications")
pdf.chapter_body(
    "The findings suggest that LightGBM provides the best stock price forecasts, "
    "allowing traders to make informed buy/sell decisions and manage risks effectively."
)

pdf.chapter_title("8. Final Deliverables")
pdf.chapter_body(
    "The project includes data preparation files, model training scripts, "
    "evaluation reports, and this presentation."
)

# Save the PDF file
pdf_filename = "Stock_Prediction_Report.pdf"
pdf.output(pdf_filename)

# Provide download links
pptx_filename, pdf_filename
